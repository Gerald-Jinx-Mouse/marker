"""Convert PPTX and XLSX files to Markdown without WeasyPrint/GTK dependencies."""

import sys
from pathlib import Path

def pptx_to_markdown(filepath: Path) -> str:
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation(str(filepath))
    lines = []

    for slide_num, slide in enumerate(prs.slides, 1):
        lines.append(f"# Slide {slide_num}")
        lines.append("")

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        # Detect likely headings (larger font or bold)
                        is_heading = False
                        if paragraph.runs:
                            run = paragraph.runs[0]
                            if run.font.bold:
                                is_heading = True
                            if run.font.size and run.font.size > Emu(1800000):  # ~14pt+
                                is_heading = True

                        if is_heading:
                            lines.append(f"## {text}")
                        else:
                            # Check if it's a bullet point
                            if paragraph.level and paragraph.level > 0:
                                indent = "  " * (paragraph.level - 1)
                                lines.append(f"{indent}- {text}")
                            elif hasattr(paragraph, '_pPr') and paragraph._pPr is not None:
                                buNone = paragraph._pPr.find(
                                    '{http://schemas.openxmlformats.org/drawingml/2006/main}buNone'
                                )
                                if buNone is None:
                                    lines.append(f"- {text}")
                                else:
                                    lines.append(text)
                            else:
                                lines.append(text)
                lines.append("")

            if shape.has_table:
                table = shape.table
                # Header row
                headers = [cell.text.strip() for cell in table.rows[0].cells]
                lines.append("| " + " | ".join(headers) + " |")
                lines.append("| " + " | ".join(["---"] * len(headers)) + " |")
                # Data rows
                for row in list(table.rows)[1:]:
                    cells = [cell.text.strip() for cell in row.cells]
                    lines.append("| " + " | ".join(cells) + " |")
                lines.append("")

        lines.append("---")
        lines.append("")

    return "\n".join(lines)


def xlsx_to_markdown(filepath: Path) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(str(filepath), data_only=True)
    lines = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        lines.append(f"# {sheet_name}")
        lines.append("")

        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            lines.append("*(empty sheet)*")
            lines.append("")
            continue

        # Find the actual data bounds (skip fully empty rows/cols)
        non_empty_rows = []
        for row in rows:
            if any(cell is not None for cell in row):
                non_empty_rows.append(row)

        if not non_empty_rows:
            lines.append("*(empty sheet)*")
            lines.append("")
            continue

        # Determine max columns with data
        max_col = 0
        for row in non_empty_rows:
            for i in range(len(row) - 1, -1, -1):
                if row[i] is not None:
                    max_col = max(max_col, i + 1)
                    break

        # Build markdown table
        for row_idx, row in enumerate(non_empty_rows):
            cells = [str(row[i]) if i < len(row) and row[i] is not None else "" for i in range(max_col)]
            lines.append("| " + " | ".join(cells) + " |")
            if row_idx == 0:
                lines.append("| " + " | ".join(["---"] * max_col) + " |")

        lines.append("")

    return "\n".join(lines)


def main():
    base_dir = Path(__file__).parent
    output_dir = base_dir / "output"
    output_dir.mkdir(exist_ok=True)

    files = [
        base_dir / "814_2_March_BSAN_760_PhantomSea.pptx",
        base_dir / "814_2_March_BSAN_760_PhantomSea.xlsx",
    ]

    for filepath in files:
        if not filepath.exists():
            print(f"File not found: {filepath}")
            continue

        print(f"Converting {filepath.name}...")

        if filepath.suffix == ".pptx":
            md = pptx_to_markdown(filepath)
        elif filepath.suffix == ".xlsx":
            md = xlsx_to_markdown(filepath)
        else:
            print(f"  Unsupported format: {filepath.suffix}")
            continue

        out_path = output_dir / f"{filepath.stem}_{filepath.suffix[1:]}.md"
        out_path.write_text(md, encoding="utf-8")
        print(f"  -> {out_path}")

    print("\nDone!")


if __name__ == "__main__":
    main()
